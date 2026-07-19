# -*- coding: utf-8 -*-
"""测试夹具：在临时目录里生成各类样本文件。"""
import os
import sys
import tempfile

# 把项目根目录(docsearch)加入 sys.path，方便 import extractor/indexer/app
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if ROOT not in sys.path:
    sys.path.insert(0, ROOT)

CJK_FONTS = [
    "/System/Library/Fonts/STHeiti Medium.ttc",
    "/System/Library/Fonts/PingFang.ttc",
    "/System/Library/Fonts/Hiragino Sans GB.ttc",
    "/Library/Fonts/Arial Unicode.ttf",
]


def make_tmpdir():
    return tempfile.mkdtemp(prefix="dstest_")


def write_txt(path, text, enc="utf-8"):
    with open(path, "w", encoding=enc) as f:
        f.write(text)


def make_docx(path, paragraphs, tables=None):
    from docx import Document
    d = Document()
    for p in paragraphs:
        d.add_paragraph(p)
    if tables:
        for t in tables:
            tb = d.add_table(rows=len(t), cols=len(t[0]))
            for i, row in enumerate(t):
                for j, val in enumerate(row):
                    tb.cell(i, j).text = val
    d.save(path)


def make_pdf(path, text):
    """生成文本型 PDF（默认字体仅含 Latin，text 限 ASCII/Latin）。"""
    import fitz
    doc = fitz.open()
    pg = doc.new_page()
    pg.insert_text((72, 72), text)
    doc.save(path)
    doc.close()


def make_pdf_cjk(path, text):
    """用系统 CJK 字体生成中文 PDF；无可用字体返回 False（测试应跳过）。"""
    import fitz
    font = next((f for f in CJK_FONTS if os.path.exists(f)), None)
    if not font:
        return False
    doc = fitz.open()
    pg = doc.new_page()
    pg.insert_text((72, 100), text, fontname="F0", fontfile=font, fontsize=14)
    doc.save(path)
    doc.close()
    return True


def make_xlsx(path, rows):
    """rows: 二维列表，写入活动 sheet。"""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row, start=1):
            ws.cell(row=r, column=c, value=val)
    wb.save(path)


def make_pptx(path, texts):
    """texts: 每个元素生成一页，作为该页文本框内容。"""
    from pptx import Presentation
    prs = Presentation()
    blank = prs.slide_layouts[6]  # 空白版式
    for t in texts:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(0, 0, 3000000, 500000)
        box.text_frame.text = t
    prs.save(path)


def make_epub(path, chapters):
    """chapters: list[str]，每段写成一章 xhtml 正文（utf-8）。"""
    import zipfile
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("mimetype", "application/epub+zip")
        for i, body in enumerate(chapters):
            xhtml = ('<?xml version="1.0" encoding="utf-8"?>\n'
                     '<html xmlns="http://www.w3.org/1999/xhtml"><body><p>%s</p></body></html>') % body
            zf.writestr("OEBPS/ch%d.xhtml" % i, xhtml)


def make_rtf(path, text):
    """造最小 rtf；为避开中文转义复杂度，正文建议用 ASCII。"""
    safe = text.replace("\\", "\\\\").replace("{", "\\{").replace("}", "\\}")
    with open(path, "w", encoding="utf-8") as f:
        f.write("{\\rtf1\\ansi " + safe + "}")


def make_odt(path, text):
    """造最小 odt：zip + content.xml（含一段文本）。"""
    import zipfile
    content_xml = (
        '<?xml version="1.0" encoding="utf-8"?>\n'
        '<office:document-content '
        'xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" '
        'xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0">'
        '<office:body><office:text><text:p>%s</text:p></office:text></office:body>'
        '</office:document-content>'
    ) % text
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("mimetype", "application/vnd.oasis.opendocument.text")
        zf.writestr("content.xml", content_xml)
