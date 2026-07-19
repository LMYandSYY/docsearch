# -*- coding: utf-8 -*-
"""文档文字提取：PDF（含扫描件 OCR）、Word(.docx/.doc)、纯文本。

所有第三方库都在函数内部延迟导入，缺少某个库时只会跳过对应格式，
不会让整个程序无法启动。
"""
import html
import os
import re
import shutil
import subprocess
import tempfile
import zipfile

# 支持的文件后缀（小写、含点）。提取逻辑按 ext 分发，见 extract_text。
SUPPORTED = {
    ".pdf", ".docx", ".doc",
    ".txt", ".md", ".markdown", ".csv", ".log",
    ".json", ".yaml", ".yml", ".xml", ".sql",
    ".js", ".py", ".cpp", ".java", ".css", ".go", ".ts",
    ".html", ".htm",
    ".xlsx", ".pptx",
    ".rtf", ".epub",
    ".odt", ".ods", ".odp",
    ".xls", ".ppt",
}

# 纯文本类：直接读文件，复用 _textfile 的多编码兜底（utf-8→gb18030→big5→latin-1）
TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".csv", ".log",
    ".json", ".yaml", ".yml", ".xml", ".sql",
    ".js", ".py", ".cpp", ".java", ".css", ".go", ".ts",
}


def extract_text(path, ocr_pages=True):
    """提取一个文件的纯文本。

    返回 (text, meta)。meta 含 ext / pages / ocr_used / errors。
    """
    ext = os.path.splitext(path)[1].lower()
    text, pages, ocr_used, errors = "", None, False, []
    try:
        if ext == ".pdf":
            text, pages, ocr_used = _pdf(path, ocr_pages)
        elif ext == ".docx":
            text = _docx_text(path)
        elif ext == ".doc":
            text = _doc_text(path)
        elif ext == ".xlsx":
            text = _xlsx_text(path)
        elif ext == ".xls":
            text = _xls_text(path)
        elif ext == ".pptx":
            text = _pptx_text(path)
        elif ext == ".ppt":
            text = _ppt_text(path)
        elif ext == ".epub":
            text = _epub_text(path)
        elif ext == ".rtf":
            text = _rtf_text(path)
        elif ext in (".odt", ".ods", ".odp"):
            text = _odf_text(path)
        elif ext in (".html", ".htm"):
            text = _strip_tags(_textfile(path))
        elif ext in TEXT_EXTS:
            text = _textfile(path)
        else:
            errors.append("不支持的格式: " + ext)
    except Exception as e:  # noqa: BLE001
        errors.append("提取失败: " + repr(e))
    return text, {"ext": ext, "pages": pages, "ocr_used": ocr_used, "errors": errors}


# ----------------------------- PDF ----------------------------- #

def _best_lang(langs):
    if "chi_sim" in langs and "eng" in langs:
        return "chi_sim+eng"
    if "chi_sim" in langs:
        return "chi_sim"
    return "eng"


def _pdf(path, ocr_pages):
    try:
        import pymupdf as fitz
    except ImportError:
        import fitz  # 旧版本别名

    doc = fitz.open(path)
    pages_text, ocr_used = [], False
    ocr_ok, lang = False, "eng"
    if ocr_pages:
        try:
            import pytesseract  # noqa: F401
            from PIL import Image  # noqa: F401
            ocr_ok = True
            try:
                lang = _best_lang(set(pytesseract.get_languages() or []))
            except Exception:
                lang = "eng"
        except Exception:
            ocr_ok = False

    for page in doc:
        t = (page.get_text() or "").strip()
        # 文字很少，基本可判定为扫描页 → OCR
        if len(t) < 10 and ocr_ok:
            try:
                import io
                import pytesseract
                from PIL import Image
                pix = page.get_pixmap(dpi=200)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                try:
                    t = (pytesseract.image_to_string(img, lang=lang) or "").strip()
                except Exception:
                    t = (pytesseract.image_to_string(img, lang="eng") or "").strip()
                ocr_used = True
            except Exception:
                pass
        pages_text.append(t)

    n = len(pages_text)
    doc.close()
    return "\n".join(pages_text), n, ocr_used


# ----------------------------- Word ----------------------------- #

def _docx_text(path):
    import docx
    d = docx.Document(path)
    parts = [p.text for p in d.paragraphs if p.text and p.text.strip()]
    for table in d.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text and cell.text.strip():
                    parts.append(cell.text)
    return "\n".join(parts)


def _doc_text(path):
    """老版 .doc：先用纯 Python(olefile) 提取正文；提取不到再退回 LibreOffice 转 .docx。"""
    text = _doc_text_olefile(path)
    if text.strip():
        return text
    docx_path = _doc_to_docx(path)
    if docx_path:
        try:
            return _docx_text(docx_path)
        except Exception:
            return ""
    return ""


def _doc_text_olefile(path):
    """纯 Python 提取 .doc 正文：读 WordDocument 流，按 UTF-16LE 解码并清理。

    不依赖任何外部程序，开箱即用；老版 Word（含中文）正文通常以 UTF-16LE 存储，
    用于关键词检索足够可靠。富文本/排版会丢失，需要精美预览可另装 LibreOffice。
    """
    try:
        import olefile
    except ImportError:
        return ""
    try:
        ole = olefile.OleFileIO(path)
        if not ole.exists("WordDocument"):
            ole.close()
            return ""
        data = ole.openstream("WordDocument").read()
        ole.close()
    except Exception:
        return ""
    return _clean_doc_text(data.decode("utf-16-le", errors="ignore"))


def _is_text_cp(o):
    """是否为「正常文字」码点：ASCII / 拉丁扩展 / CJK / 全角 / 通用标点。"""
    return (
        0x20 <= o <= 0x7E
        or 0xA0 <= o <= 0x24F
        or 0x3000 <= o <= 0x9FFF      # CJK 符号/统一表意/扩展A
        or 0xF900 <= o <= 0xFAFF      # CJK 兼容表意
        or 0xFF00 <= o <= 0xFFEF      # 全角形式
        or 0x2000 <= o <= 0x206F      # 通用标点（破折号、引号等）
    )


def _clean_doc_text(text):
    """把 .doc 二进制头解码产生的噪声清掉，只留可读正文。

    1) 非文字码点统一转成空格（而非丢弃），避免不相邻的字被错误拼到一起；
    2) 正文前的 FIB 二进制头里，可读字符都是零散短串，真正正文是连续长串——
       所以定位到第一段连续长文字（≥6）作为正文起点，跳过前面的噪声。
    """
    out = []
    for ch in text:
        o = ord(ch)
        if ch in "\r\n":
            out.append("\n")
        elif ch == "\t":
            out.append(" ")
        elif _is_text_cp(o):
            out.append(ch)
        else:
            out.append(" ")
    s = "".join(out)
    s = re.sub(r"[ \t]+", " ", s)
    m = re.search(r"\S{6,}", s)
    if m:
        s = s[m.start():]
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


def _soffice_convert(path, to_ext):
    """用 LibreOffice 把任意文档转成指定格式（docx/xlsx/pptx…），返回输出路径或 None。

    供 .doc/.xls/.ppt 等老格式兜底：先转成对应新格式，再用纯 Python 库读取。
    """
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return None
    outdir = tempfile.mkdtemp(prefix="docconv_")
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", to_ext, "--outdir", outdir, path],
            check=False, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=180,
        )
    except Exception:
        return None
    out = os.path.join(outdir, os.path.splitext(os.path.basename(path))[0] + "." + to_ext)
    return out if os.path.isfile(out) else None


def _doc_to_docx(path):
    return _soffice_convert(path, "docx")


# ----------------------------- 纯文本 ----------------------------- #

def _textfile(path):
    for enc in ("utf-8", "gb18030", "big5", "latin-1"):
        try:
            with open(path, "r", encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    return ""


# ----------------------------- 表格 / 演示 ----------------------------- #

def _xlsx_text(path):
    """.xlsx：读所有 sheet 的单元格值；data_only=True 让公式取计算结果。"""
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    try:
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for v in row:
                    if v is None:
                        continue
                    s = str(v).strip()
                    if s:
                        parts.append(s)
    finally:
        wb.close()
    return "\n".join(parts)


def _xls_text(path):
    """老 .xls：优先 LibreOffice 转 xlsx 再读；转换不可用返回空串。"""
    xlsx_path = _soffice_convert(path, "xlsx")
    if xlsx_path:
        try:
            return _xlsx_text(xlsx_path)
        except Exception:
            return ""
    return ""


def _pptx_text(path):
    """.pptx：读所有幻灯片里文本框的文字。"""
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                t = shape.text_frame.text
                if t and t.strip():
                    parts.append(t)
    return "\n".join(parts)


def _ppt_text(path):
    """老 .ppt：优先 LibreOffice 转 pptx 再读；转换不可用返回空串。"""
    pptx_path = _soffice_convert(path, "pptx")
    if pptx_path:
        try:
            return _pptx_text(pptx_path)
        except Exception:
            return ""
    return ""


# ----------------------------- 富文本 / 电子书 / ODF ----------------------------- #

_TAG_RE = re.compile(r"<[^>]+>")


def _strip_tags(text):
    """粗略去 HTML/XML 标签并折叠多余空白，便于全文检索（不追求还原排版）。"""
    text = _TAG_RE.sub(" ", text)
    text = text.replace("&nbsp;", " ")
    return re.sub(r"\s+", " ", text).strip()


def _rtf_text(path):
    """.rtf：用 striprtf 转纯文本（rtf 本身是文本，可直接走多编码读取）。"""
    from striprtf.striprtf import rtf_to_text
    return rtf_to_text(_textfile(path)).strip()


def _zip_read_texts(path, suffixes):
    """从 zip 里读出所有名字以给定后缀结尾的条目，按 utf-8(容错) 解码并拼接。"""
    parts = []
    with zipfile.ZipFile(path) as zf:
        for name in zf.namelist():
            low = name.lower()
            if any(low.endswith(s) for s in suffixes):
                try:
                    parts.append(zf.read(name).decode("utf-8", errors="ignore"))
                except Exception:
                    continue
    return "\n".join(parts)


def _epub_text(path):
    """.epub：本质 zip + xhtml，读出各章节正文并去标签。"""
    return _strip_tags(_zip_read_texts(path, (".xhtml", ".html", ".htm")))


def _odf_text(path):
    """.odt/.ods/.odp：zip 里的 content.xml，去标签取正文。"""
    try:
        with zipfile.ZipFile(path) as zf:
            raw = zf.read("content.xml").decode("utf-8", errors="ignore")
    except Exception:
        return ""
    return _strip_tags(raw)


# ----------------------------- 预览 ----------------------------- #

def docx_to_html(path):
    import mammoth
    with open(path, "rb") as f:
        result = mammoth.convert_to_html(f)
    return result.value or "<p>(文档为空)</p>"


def to_preview_html(path):
    """供预览页使用的 HTML 正文。

    docx/doc 给富文本预览；其余支持格式（xlsx/pptx/epub/rtf/odt/代码等）
    一律回退为提取出的纯文本 <pre>，同样可高亮命中词。
    """
    ext = os.path.splitext(path)[1].lower()
    if ext == ".docx":
        return docx_to_html(path)
    if ext == ".doc":
        # 装了 LibreOffice 就给富文本预览；否则回退纯文本（同样可高亮命中词）
        docx_path = _doc_to_docx(path)
        if docx_path:
            try:
                return docx_to_html(docx_path)
            except Exception:
                pass
        text = _doc_text_olefile(path)
        if text.strip():
            return "<pre style='white-space:pre-wrap;word-break:break-word'>" + html.escape(text) + "</pre>"
        return "<p>无法读取此 .doc 文件的内容，可尝试另存为 .docx。</p>"
    # 其余格式：提取纯文本后 <pre> 预览
    text, _ = extract_text(path)
    if text.strip():
        return "<pre style='white-space:pre-wrap;word-break:break-word'>" + html.escape(text) + "</pre>"
    return "<p>无法提取该文件的文本内容。</p>"
